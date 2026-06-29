#!/usr/bin/env python3
import argparse
import csv
import json
import re
import sys
from pathlib import Path
from xml.etree import ElementTree

try:
    import docx
except Exception:
    docx = None

try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import pptx
except Exception:
    pptx = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    from lxml import html as lxml_html
except Exception:
    lxml_html = None


IMAGE_EXPECTATIONS = {
    "chart-page.jpg": {
        "snippets": [
            "Alphabet Inc. Class A",
            "NASDAQ Composite",
            "12/24",
        ],
        "visual_asset_min": 1,
        "note": "Chart audit verifies searchable labels, not digitized line values.",
    },
    "diagram-notes.jpg": {
        "snippets": [
            "Parser Smoke Diagram Notes",
            "Parse",
            "Chunk",
            "Index",
            "Visual parser should capture diagram labels and arrows.",
        ],
        "visual_asset_min": 1,
    },
    "large-chart-table.png": {
        "snippets": [
            "Parser Smoke Chart Data",
            "Metric",
            "Revenue",
            "Cost",
            "Margin",
        ],
        "table_cells": [
            "Metric",
            "Q1",
            "Q2",
            "Revenue",
            "100",
            "125",
            "Cost",
            "75",
            "82",
            "Margin",
            "25",
            "43",
        ],
        "table_count_min": 1,
        "visual_asset_min": 1,
    },
    "large-table.png": {
        "snippets": [
            "Parser Smoke Revenue Table",
            "Quarter",
            "Revenue",
            "Growth",
        ],
        "table_cells": [
            "Quarter",
            "Revenue",
            "Growth",
            "Q1",
            "$100M",
            "10%",
            "Q2",
            "$125M",
            "25%",
            "Q3",
            "$140M",
            "12%",
        ],
        "table_count_min": 1,
        "visual_asset_min": 1,
    },
    "mini-table.png": {
        "visual_asset_min": 1,
        "note": "Mini image fixture is only used to prove image asset handling.",
    },
}


def main():
    args = parse_args()
    parsed_docs = load_json(Path(args.parsed_docs))
    source_config = load_sources(Path(args.sources))
    report_dir = Path(args.report_dir)
    report_dir.mkdir(parents=True, exist_ok=True)

    documents_by_path = {
        document_path(document): document for document in parsed_docs.get("documents", [])
    }
    checks = []
    files = []

    for source in source_config:
        root_dir = resolve_root(Path(args.sources), source)
        for relative in source.get("files", []):
            relative_path = normalize_path(relative)
            file_path = root_dir / relative_path
            document = documents_by_path.get(relative_path)
            result = audit_file(relative_path, file_path, document)
            files.append(result)
            checks.extend(result["checks"])

    summary = summarize(checks, files)
    report = {
        "status": "passed" if summary["failed"] == 0 and summary["error"] == 0 else "failed",
        "summary": summary,
        "files": files,
    }

    (report_dir / "extraction-audit.json").write_text(json.dumps(report, indent=2), "utf-8")
    (report_dir / "extraction-audit.md").write_text(render_markdown(report), "utf-8")
    print(json.dumps({"status": report["status"], "summary": summary}))
    return 0 if report["status"] == "passed" else 1


def audit_file(relative_path, file_path, document):
    checks = []
    if document is None:
        checks.append(fail("document_exists", "Parsed document was not emitted."))
        return file_result(relative_path, checks)

    checks.append(pass_check("document_exists", "Parsed document was emitted."))
    checks.append(
        bool_check(
            "body_nonempty",
            bool(document.get("body", "").strip()),
            "Parsed body is nonempty.",
            "Parsed body is empty.",
        )
    )

    try:
        expected = expected_for_file(file_path, relative_path)
    except Exception as error:
        checks.append(fail("original_readable", f"Could not read original file: {error}"))
        return file_result(relative_path, checks)

    checks.append(pass_check("original_readable", "Original file was read independently."))
    for note in expected.get("notes", []):
        checks.append(pass_check("note", note))

    haystack = combined_search_text(document)
    for snippet in unique(expected.get("snippets", [])):
        if not material_text(snippet):
            continue
        checks.append(
            bool_check(
                "snippet_preserved",
                contains_text(haystack, snippet),
                f"Preserved snippet: {snippet}",
                f"Missing snippet: {snippet}",
            )
        )

    actual_cells = table_cell_text(document)
    cell_haystack = "\n".join([haystack, "\n".join(actual_cells)])
    for cell in unique(expected.get("table_cells", [])):
        if not material_cell(cell):
            continue
        checks.append(
            bool_check(
                "table_cell_preserved",
                contains_text(cell_haystack, cell),
                f"Preserved table cell: {cell}",
                f"Missing table cell: {cell}",
            )
        )

    table_count_min = expected.get("table_count_min", 0)
    if table_count_min > 0:
        actual_table_count = len(document.get("layout", {}).get("tables", []) or [])
        checks.append(
            bool_check(
                "table_count",
                actual_table_count >= table_count_min,
                f"Parsed {actual_table_count} table(s), expected at least {table_count_min}.",
                f"Parsed {actual_table_count} table(s), expected at least {table_count_min}.",
            )
        )

    visual_asset_min = expected.get("visual_asset_min", 0)
    if visual_asset_min > 0:
        actual_visual_count = len(document.get("layout", {}).get("visualAssets", []) or [])
        checks.append(
            bool_check(
                "visual_asset_count",
                actual_visual_count >= visual_asset_min,
                f"Parsed {actual_visual_count} visual asset(s), expected at least {visual_asset_min}.",
                f"Parsed {actual_visual_count} visual asset(s), expected at least {visual_asset_min}.",
            )
        )

    page_count_min = expected.get("page_count_min", 0)
    if page_count_min > 0:
        actual_page_count = len(document.get("layout", {}).get("pages", []) or [])
        checks.append(
            bool_check(
                "page_count",
                actual_page_count >= page_count_min,
                f"Parsed {actual_page_count} page(s), expected at least {page_count_min}.",
                f"Parsed {actual_page_count} page(s), expected at least {page_count_min}.",
            )
        )

    return file_result(relative_path, checks)


def expected_for_file(file_path, relative_path):
    suffix = file_path.suffix.lower()
    if suffix in {".txt"}:
        return text_expectations(file_path)
    if suffix in {".md", ".markdown"}:
        return markdown_expectations(file_path)
    if suffix == ".csv":
        return delimited_expectations(file_path, ",")
    if suffix == ".tsv":
        return delimited_expectations(file_path, "\t")
    if suffix == ".json":
        return json_expectations(file_path)
    if suffix == ".xml":
        return xml_expectations(file_path)
    if suffix in {".htm", ".html"}:
        return html_expectations(file_path)
    if suffix == ".docx":
        return docx_expectations(file_path)
    if suffix == ".pptx":
        return pptx_expectations(file_path)
    if suffix in {".xlsx", ".xlsm"}:
        return xlsx_expectations(file_path)
    if suffix == ".pdf":
        return pdf_expectations(file_path)
    if suffix in {".png", ".jpg", ".jpeg"}:
        return image_expectations(relative_path)
    return {"snippets": meaningful_lines(file_path.read_text("utf-8", errors="replace"))[:20]}


def text_expectations(file_path):
    return {"snippets": meaningful_lines(file_path.read_text("utf-8", errors="replace"))}


def markdown_expectations(file_path):
    text = file_path.read_text("utf-8", errors="replace")
    snippets = meaningful_lines(text)
    table_cells, table_count = markdown_table_cells(text)
    return {
        "snippets": snippets,
        "table_cells": table_cells,
        "table_count_min": table_count,
    }


def delimited_expectations(file_path, delimiter):
    rows = []
    with file_path.open("r", encoding="utf-8", newline="") as handle:
        for row in csv.reader(handle, delimiter=delimiter):
            rows.append(row)
    cells = [cell for row in rows for cell in row]
    return {
        "snippets": rows[0] if rows else [],
        "table_cells": cells,
        "table_count_min": 1 if rows else 0,
    }


def json_expectations(file_path):
    data = json.loads(file_path.read_text("utf-8"))
    snippets = flatten_json_values(data)
    return {"snippets": snippets[:80]}


def xml_expectations(file_path):
    root = ElementTree.parse(file_path).getroot()
    snippets = []
    for element in root.iter():
        tag = element.tag.split("}")[-1]
        if material_text(tag):
            snippets.append(tag)
        if element.text and material_text(element.text):
            snippets.append(element.text)
        for value in element.attrib.values():
            if material_text(value):
                snippets.append(value)
    return {"snippets": snippets[:80]}


def html_expectations(file_path):
    raw = file_path.read_text("utf-8", errors="replace")
    raw = sec_text_body(raw)
    if lxml_html is None:
        return {"snippets": meaningful_lines(strip_html(raw))[:40]}
    document = lxml_html.fromstring(raw.encode("utf-8"))
    for bad in document.xpath("//script|//style"):
        bad.drop_tree()
    drop_nonvisible_html(document)
    snippets = []
    for text in document.xpath("//body//text() | //text()[not(ancestor::script) and not(ancestor::style)]"):
        snippets.extend(meaningful_lines(str(text)))
    tables = material_html_tables(document)
    cells = [cell.text_content().strip() for table in tables for cell in table.xpath(".//th|.//td")]
    return {
        "snippets": snippets[:80],
        "table_cells": cells,
        "table_count_min": 1 if tables else 0,
    }


def docx_expectations(file_path):
    if docx is None:
        raise RuntimeError("python-docx is not installed")
    document = docx.Document(str(file_path))
    snippets = [paragraph.text for paragraph in document.paragraphs if material_text(paragraph.text)]
    cells = []
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                cells.append(cell.text)
    return {
        "snippets": snippets + cells,
        "table_cells": cells,
        "table_count_min": len(document.tables),
    }


def pptx_expectations(file_path):
    if pptx is None:
        raise RuntimeError("python-pptx is not installed")
    presentation = pptx.Presentation(str(file_path))
    snippets = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and material_text(shape.text):
                snippets.extend(meaningful_lines(shape.text))
    return {"snippets": snippets}


def xlsx_expectations(file_path):
    if openpyxl is None:
        raise RuntimeError("openpyxl is not installed")
    workbook = openpyxl.load_workbook(str(file_path), data_only=False)
    snippets = []
    visual_count = 0
    table_count = 0
    for sheet in workbook.worksheets:
        snippets.append(sheet.title)
        visual_count += len(getattr(sheet, "_charts", []) or [])
        visual_count += len(getattr(sheet, "_images", []) or [])
        sheet_has_rows = False
        for row in sheet.iter_rows():
            row_has_values = False
            for cell in row:
                value = cell.value
                if value is None:
                    continue
                row_has_values = True
                value_text = str(value)
                if value_text.startswith("=") or not numeric_like(value_text):
                    snippets.append(value_text)
            sheet_has_rows = sheet_has_rows or row_has_values
        if sheet_has_rows:
            table_count += 1
    return {
        "snippets": snippets[:120],
        "table_count_min": table_count,
        "visual_asset_min": visual_count,
    }


def pdf_expectations(file_path):
    if pdfplumber is None:
        raise RuntimeError("pdfplumber is not installed")
    snippets = []
    table_cells = []
    page_count = 0
    with pdfplumber.open(str(file_path)) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            text = page.extract_text() or ""
            snippets.extend(meaningful_lines(text)[:20])
            cells, _ = markdown_table_cells(text)
            table_cells.extend(cells)
            for table in page.extract_tables() or []:
                for row in table:
                    for cell in row:
                        if cell:
                            table_cells.append(cell)
    return {
        "snippets": snippets[:80],
        "table_cells": table_cells,
        "table_count_min": 1 if table_cells else 0,
        "page_count_min": page_count,
    }


def image_expectations(relative_path):
    expectation = dict(IMAGE_EXPECTATIONS.get(relative_path, {"visual_asset_min": 1}))
    notes = []
    if expectation.get("note"):
        notes.append(expectation.pop("note"))
    expectation["notes"] = notes
    return expectation


def combined_search_text(document):
    parts = [document.get("title", ""), document.get("body", "")]
    metadata = document.get("metadata", {}) or {}
    for value in metadata.values():
        if isinstance(value, str):
            parts.append(value)
    parts.extend(table_cell_text(document))
    for region in document.get("layout", {}).get("regions", []) or []:
        text = region.get("text")
        if isinstance(text, str):
            parts.append(text)
    return "\n".join(parts)


def table_cell_text(document):
    cells = []
    for table in document.get("layout", {}).get("tables", []) or []:
        for cell in table.get("cells", []) or []:
            text = cell.get("text")
            if isinstance(text, str):
                cells.append(text)
    return cells


def markdown_table_cells(text):
    cells = []
    table_count = 0
    in_table = False
    for line in text.splitlines():
        if "|" not in line:
            in_table = False
            continue
        parts = [part.strip() for part in line.strip().strip("|").split("|")]
        if len(parts) < 2:
            continue
        if all(re.fullmatch(r":?-{2,}:?", part.replace(" ", "")) for part in parts):
            continue
        if not in_table:
            table_count += 1
            in_table = True
        cells.extend(parts)
    return cells, table_count


def flatten_json_values(value):
    flattened = []
    if isinstance(value, dict):
        for key, child in value.items():
            flattened.append(str(key))
            flattened.extend(flatten_json_values(child))
    elif isinstance(value, list):
        for item in value:
            flattened.extend(flatten_json_values(item))
    elif value is not None:
        flattened.append(str(value))
    return [item for item in flattened if material_text(item)]


def meaningful_lines(text):
    lines = []
    for line in text.splitlines():
        clean = line.strip()
        if material_text(clean):
            lines.append(clean)
    return lines[:120]


def strip_html(raw):
    return re.sub(r"<[^>]+>", " ", raw)


def drop_nonvisible_html(document):
    hidden_style = "ABCDEFGHIJKLMNOPQRSTUVWXYZ"
    visible_style = "abcdefghijklmnopqrstuvwxyz"
    for node in list(
        document.xpath(
            "//*[contains(translate(@style, "
            f"'{hidden_style}', '{visible_style}'), 'display:none')]"
        )
    ):
        node.drop_tree()
    for node in list(document.xpath("//*[@hidden]")):
        node.drop_tree()


def material_html_tables(document):
    tables = []
    for table in document.xpath("//table"):
        text = table.text_content()
        if material_text(text):
            tables.append(table)
    return tables


def sec_text_body(raw):
    match = re.search(r"<TEXT>(.*)</TEXT>", raw, flags=re.IGNORECASE | re.DOTALL)
    return match.group(1) if match else raw


def material_text(value):
    return len(normalize_alnum(value)) >= 3


def material_cell(value):
    normalized = normalize_alnum(value)
    return len(normalized) >= 2 or bool(re.search(r"\d", str(value)))


def contains_text(haystack, needle):
    haystack_normalized = normalize_text(haystack)
    needle_normalized = normalize_text(needle)
    if needle_normalized and needle_normalized in haystack_normalized:
        return True
    haystack_alnum = normalize_alnum(haystack)
    needle_alnum = normalize_alnum(needle)
    if needle_alnum and needle_alnum in haystack_alnum:
        return True
    haystack_compact = normalize_compact(haystack)
    needle_compact = normalize_compact(needle)
    if needle_compact and needle_compact in haystack_compact:
        return True
    return material_token_coverage(haystack, needle) >= 0.85


def normalize_text(value):
    return re.sub(r"\s+", " ", str(value).lower()).strip()


def normalize_alnum(value):
    return re.sub(r"[^a-z0-9]+", " ", str(value).lower()).strip()


def normalize_compact(value):
    return re.sub(r"[^a-z0-9]+", "", str(value).lower())


def material_token_coverage(haystack, needle):
    tokens = material_tokens(needle)
    if len(tokens) < 4:
        return 0.0
    haystack_tokens = set(material_tokens(haystack))
    if not haystack_tokens:
        return 0.0
    matched = sum(1 for token in tokens if token in haystack_tokens)
    return matched / len(tokens)


def material_tokens(value):
    tokens = re.findall(r"[a-z0-9]+", str(value).lower())
    material = []
    for token in tokens:
        if token.isdigit() or len(token) >= 3:
            material.append(token)
    return unique(material)


def numeric_like(value):
    return bool(re.fullmatch(r"[$€£]?-?\d+(?:[.,]\d+)?%?", str(value).strip()))


def document_path(document):
    provenance = document.get("provenance", {}) or {}
    metadata = document.get("metadata", {}) or {}
    return normalize_path(provenance.get("path") or metadata.get("relativePath") or document.get("id", ""))


def normalize_path(value):
    return str(value).replace("\\", "/")


def unique(values):
    seen = set()
    out = []
    for value in values:
        key = normalize_alnum(value)
        if not key or key in seen:
            continue
        seen.add(key)
        out.append(str(value).strip())
    return out


def pass_check(name, message):
    return {"name": name, "status": "passed", "message": message}


def fail(name, message):
    return {"name": name, "status": "failed", "message": message}


def bool_check(name, condition, pass_message, fail_message):
    return pass_check(name, pass_message) if condition else fail(name, fail_message)


def file_result(relative_path, checks):
    failed = any(check["status"] == "failed" for check in checks)
    error = any(check["status"] == "error" for check in checks)
    status = "error" if error else ("failed" if failed else "passed")
    return {
        "path": relative_path,
        "status": status,
        "checkCount": len(checks),
        "failed": sum(1 for check in checks if check["status"] == "failed"),
        "checks": checks,
    }


def summarize(checks, files):
    return {
        "fileCount": len(files),
        "passedFiles": sum(1 for file in files if file["status"] == "passed"),
        "failedFiles": sum(1 for file in files if file["status"] == "failed"),
        "errorFiles": sum(1 for file in files if file["status"] == "error"),
        "checkCount": len(checks),
        "passed": sum(1 for check in checks if check["status"] == "passed"),
        "failed": sum(1 for check in checks if check["status"] == "failed"),
        "error": sum(1 for check in checks if check["status"] == "error"),
    }


def render_markdown(report):
    lines = [
        "# Original Extraction Audit",
        "",
        f"- Status: {report['status']}",
        f"- Files: {report['summary']['fileCount']}",
        f"- Checks: {report['summary']['checkCount']}",
        f"- Failed checks: {report['summary']['failed']}",
        "",
        "| File | Status | Checks | Failed |",
        "| --- | --- | ---: | ---: |",
    ]
    for file in report["files"]:
        lines.append(
            f"| `{file['path']}` | {file['status']} | {file['checkCount']} | {file['failed']} |"
        )
    failed_files = [file for file in report["files"] if file["status"] != "passed"]
    if failed_files:
        lines.extend(["", "## Failures", ""])
        for file in failed_files:
            lines.append(f"### {file['path']}")
            for check in file["checks"]:
                if check["status"] != "passed":
                    lines.append(f"- {check['name']}: {check['message']}")
            lines.append("")
    return "\n".join(lines)


def load_json(path):
    return json.loads(path.read_text("utf-8"))


def load_sources(path):
    parsed = load_json(path)
    if isinstance(parsed, list):
        return parsed
    sources = parsed.get("sources")
    if not isinstance(sources, list):
        raise ValueError(f"{path} must be an array or an object with sources[]")
    return sources


def resolve_root(sources_path, source):
    root_dir = Path(source.get("rootDir", "."))
    if root_dir.is_absolute():
        return root_dir
    return (sources_path.parent / root_dir).resolve()


def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument("--parsed-docs", required=True)
    parser.add_argument("--sources", required=True)
    parser.add_argument("--report-dir", required=True)
    return parser.parse_args()


if __name__ == "__main__":
    sys.exit(main())
